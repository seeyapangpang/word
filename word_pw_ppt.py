import json
import streamlit as st
import pandas as pd
import time
from io import BytesIO
from openai import OpenAI
import requests
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches

# ✅ 비밀번호 보호 기능
def check_password():
    if "password_correct" not in st.session_state:
        st.session_state["password_correct"] = False

    if not st.session_state["password_correct"]:
        password = st.text_input("비밀번호를 입력하세요:", type="password")
        if password == st.secrets["APP_PASSWORD"]:  
            st.session_state["password_correct"] = True
        else:
            st.error("비밀번호가 틀렸습니다.")
            return False
    return True

# ✅ 실시간 환율 가져오기
def get_exchange_rate():
    try:
        response = requests.get("https://api.exchangerate-api.com/v4/latest/USD")
        data = response.json()
        return data["rates"].get("KRW", 1300)
    except Exception:
        return 1300

# ✅ 예상 비용 및 시간 계산 함수
def estimate_cost(word_count, avg_example_length=50):
    token_per_word = 2  
    token_per_example = avg_example_length * 1.2  
    total_tokens = word_count * (token_per_word + token_per_example)
    cost_per_1k_tokens = 0.0015  
    usd_cost = (total_tokens / 1000) * cost_per_1k_tokens
    exchange_rate = get_exchange_rate()
    krw_cost = usd_cost * exchange_rate
    estimated_time = word_count * 0.2  
    return total_tokens, usd_cost, krw_cost, exchange_rate, estimated_time

# ✅ 번역 및 예문 생성 함수 (배치 크기 조정 및 오류 처리 추가)
def generate_batch_translations(words, client, retries=3):
    for attempt in range(retries):
        try:
            words_string = json.dumps(words)  
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant. Always respond in the following JSON format: "
                                                     '{"translations": [{"word": "<word>", "ipa": "<IPA pronunciation>", "korean": "<korean translations>", "example": "<short English sentence>", "example_korean": "<Korean translation of the example>"}]}'},
                    {"role": "user", "content": f"Provide IPA pronunciation, a list of Korean translations, and a very short English sentence for toddlers along with its Korean translation. Here are the words: {words_string}."}
                ]
            )
            output = response.choices[0].message.content.strip()
            parsed_response = json.loads(output)
            translations = []
            for item in parsed_response.get("translations", []):
                word = item.get("word", "").strip()
                ipa = item.get("ipa", "발음 없음").strip()
                ipa = f"[ {ipa.replace('/', '').strip()} ]" if ipa != "발음 없음" else "발음 없음"
                korean = item.get("korean", "번역 없음").strip()
                example = item.get("example", "No example available").strip()
                example_korean = item.get("example_korean", "예문 없음").strip()
                combined_example = f"{example} ({example_korean})"
                translations.append([word, ipa, korean, combined_example, example, example_korean])
            return translations
        except Exception as e:
            if attempt < retries - 1:
                time.sleep(2)  # 재시도 전 대기
                continue
            else:
                return [[word, "발음 없음", "번역 없음", "예문 오류", "예문 오류", "예문 오류"] for word in words]

# ✅ 배치 크기 유지 (10)
batch_size = 10

# ✅ 엑셀 생성 함수
def write_to_excel(result_df):
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        result_df.to_excel(writer, index=False)
    return output.getvalue()

# ✅ 파워포인트 생성 함수
def write_to_pptx(result_df):
    prs = Presentation()
    for _, row in result_df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = row['Word']
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        text_frame = textbox.text_frame
        text_frame.text = f"IPA: {row['IPA']}\n\nKorean: {row['Korean']}\n\nExample: {row['Combined Example']}"
    output = BytesIO()
    prs.save(output)
    return output.getvalue()

# ✅ 비밀번호 확인 후 실행
if check_password():
    st.title("단어 번역 및 예문 생성기")
    st.write("엑셀 파일을 업로드하면 단어에 대한 IPA 발음, 번역, 예문을 자동 생성합니다.")

    client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

    uploaded_file = st.file_uploader("엑셀 파일을 업로드하세요", type=["xlsx"])

    if "result_df" not in st.session_state:
        st.session_state.result_df = None

    if uploaded_file is not None:
        df = pd.read_excel(uploaded_file, header=None)
        df = df.iloc[1:, :1]
        df.columns = ["Word"]

        word_count = len(df)
        total_tokens, usd_cost, krw_cost, exchange_rate, estimated_time = estimate_cost(word_count)

        st.write("업로드된 데이터:")
        st.write(df)

        st.subheader("예상 비용 및 시간")
        st.write(f"- 예상 토큰 수: {total_tokens}")
        st.write(f"- 예상 비용 (USD): ${usd_cost:.4f}")
        st.write(f"- 예상 비용 (KRW): {krw_cost:,.0f}원 (환율: {exchange_rate:.2f} KRW/USD)")
        st.write(f"- 예상 시간: {estimated_time:.2f} 초")

        if st.button("Go (API 요청 시작)"):
            start_time = time.time()
            st.write("번역과 예문을 생성하는 중입니다...")

            translations = []
            progress_bar = st.progress(0)

            for i in range(0, word_count, batch_size):
                batch_words = df["Word"].iloc[i:i + batch_size].tolist()
                translations.extend(generate_batch_translations(batch_words, client))
                progress_bar.progress(min((i + batch_size) / word_count, 1.0))

            end_time = time.time()
            execution_time = end_time - start_time
            st.write(f"실제 소요 시간: {execution_time:.2f} 초")
            
            st.session_state.result_df = pd.DataFrame(translations, columns=["Word", "IPA", "Korean", "Combined Example", "English Example", "Korean Example"])

    if st.session_state.result_df is not None:
        st.subheader("번역 및 예문 생성 결과")
        st.write(st.session_state.result_df)
        st.download_button(
            label="결과 다운로드 (엑셀)",
            data=write_to_excel(st.session_state.result_df),
            file_name="translated_vocabulary.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        st.download_button(
            label="결과 다운로드 (파워포인트)",
            data=write_to_pptx(st.session_state.result_df),
            file_name="translated_vocabulary.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

이 코딩을 그대로 유지하려고해. 그런데 api를 받는 과정에서 데이터를 못 받는 경우도 생기더라구. 이를 ㄹ방지하기 위해 3번까지 요청하는 것으로 해줘 
